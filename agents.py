from langchain_community.document_loaders import YoutubeLoader,WebBaseLoader,UnstructuredWordDocumentLoader,UnstructuredPowerPointLoader,AzureAIDocumentIntelligenceLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_chroma import Chroma
# from langchain_community.vectorstores import chroma
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings.sentence_transformer import (
    SentenceTransformerEmbeddings,
)
from langchain.prompts import ChatPromptTemplate,PromptTemplate
from langchain.llms import HuggingFaceHub
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
# from langchain.llms import HuggingFaceHub
from langchain_community.llms import HuggingFaceHub
from langchain.chains.summarize import load_summarize_chain
import streamlit as st
import re
# import chromadb

def is_valid_youtube_url(url):
    # Regular expression to match YouTube video URLs
    youtube_regex = (
        r'(https?://)?(www\.)?'
        '(youtube|youtu|youtube-nocookie)\.(com|be)/'
        '(watch\?v=|embed/|v/|.+\?v=)?([^&=%\?]{11})')

    # Match the provided URL with the regex pattern
    match = re.match(youtube_regex, url)
    
    # If there's a match, it's a valid YouTube URL
    if match:
        return True
    else:
        return False
# --------------------------------------------------- website --------------------------------------------------------

def load_web_documents(url):
    # get the text in document form
    loader = WebBaseLoader(url)
    data = loader.load()
    return data

# --------------------------------------------------- you tube ------------------------------------------------------------

def load_url_documents(link):

    if is_valid_youtube_url(link)==False:
        return None
    loader = YoutubeLoader.from_youtube_url(
        link,
        add_video_info=True,
        language=["en","hi","id"],
        translation="en",
    )
    documents = loader.load()
    return documents
# ----------------------------------------------------------------------- pdf -----------------------------------
from PyPDF2 import PdfReader
from langchain_community.document_loaders import PyMuPDFLoader
import tempfile

def save_uploadedfile(uploadedfile):
    with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
        tmp_file.write(uploadedfile.read())
        return tmp_file.name

def load_pdf_documents(pdf_docs):
    data=[]
    for pdf in pdf_docs:
        path=save_uploadedfile(pdf)
        loader = PyMuPDFLoader(path)
        data.append(loader.load())
    return data[0]

# --------------------------------------------------------- DOC --------------------------------------
from langchain_community.document_loaders import Docx2txtLoader
def load_doc_documents(pdf_docs):
    # loader = UnstructuredWordDocumentLoader(
    #     pdf_docs, mode="elements", strategy="fast",

    # ) 
    # data=loader.load()
    # return data

    data=[]
    for pdf in pdf_docs:
        path=save_uploadedfile(pdf)
        loader = Docx2txtLoader(path)
        data.append(loader.load())
    return data[0]

# ---------------------------------------------------------- PPT ----------------------------------------------------------

# pip install --upgrade --quiet  langchain langchain-community azure-ai-documentintelligence
from langchain.schema import Document 
from pptx import Presentation
def load_ppt_documents(pdf_docs):
    # loader = UnstructuredPowerPointLoader(
    #     pdf_docs, mode="elements", strategy="fast",

    # ) 
    # data=loader.load()
    # return data

    # data=[]
    # for file in pdf_docs:
    #     path=save_uploadedfile(file)
    #     loader = UnstructuredPowerPointLoader(path) 
    #     data.append(loader.load())
    # return data[0]
    text=""
    for file in pdf_docs:
        try:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            st.error("Error: " + str(e))

    # Create a Document object with page_content attribute
    document = Document(page_content=text)
    return [document]

# # -------------------------------------- CSV ------------------------------------------------------------------------

# from langchain_community.document_loaders.csv_loader import CSVLoader

# def load_csv_documents(csv_files):
#     data=[]
#     for file in csv_files:
#         path=save_uploadedfile(file)
#         loader = CSVLoader(path)
#         data.append(loader.load())
#     return data[0]

# ---------------------------------------------------------------------------------------------------------------------

def load_documents(url,pdf_docs,file_type):
    if file_type == "PDF":
        return load_pdf_documents(pdf_docs)
    elif file_type=="DOC":
        return load_doc_documents(pdf_docs)
    elif file_type=="PPT":
        return load_ppt_documents(pdf_docs)
    elif file_type == "YOUTUBE":
        return load_url_documents(url)
    elif file_type=="WEBSITE":
        return load_web_documents(url)
    else:
        return ""


def split_documents(documents):
    # st.write(documents)
    
    text_splitter = RecursiveCharacterTextSplitter(
        separators=[
            "\n\n",
            "\n",
            " ",
            ".",
            ",",
            "",
        ],
        chunk_size=1000,
        chunk_overlap=50,
        length_function=len,
        is_separator_regex=False,
    )

    # Make splits
    splits = text_splitter.split_documents(documents)
    return splits

# create the open-source embedding function
embedding_function = SentenceTransformerEmbeddings(model_name="all-MiniLM-L6-v2")
def define_vectore_store(splits):

    # load it into FAISS
    db = FAISS.from_documents(documents=splits, embedding=embedding_function)
    return db

def retriever_chain(db):
    retriever = db.as_retriever()

    # Prompt
    template = """Answer the question based only on the following context:
    {context}

    Question: {question}
    """

    prompt = ChatPromptTemplate.from_template(template)

    # llm
    HUGGINGFACEHUB_API_TOKEN = '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'
    repo_id = "mistralai/Mistral-7B-Instruct-v0.2"
    # repo_id="meta-llama/Meta-Llama-3-8B"
    llm = HuggingFaceHub(repo_id=repo_id,model_kwargs={"temperature":0.5, "max_length":1024},huggingfacehub_api_token=HUGGINGFACEHUB_API_TOKEN)

    rag_chain = (
        {"context": retriever, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    return rag_chain

def main():
    st.set_page_config(page_title="Chat with your Data",
                   page_icon='🤖',)
    st.title('Chat with your Data 💁')
    st.markdown('Supports 5 types of data format PDF/DOC/PPT/YOUTUBE/WEBSITE')
    
    file_type = st.selectbox("Select the type of data:", ["PDF","YOUTUBE","WEBSITE","DOC","PPT"])


    # Sidebar for PDF input
    st.sidebar.header("Browse document (PDF/DOC/PPT)")
    pdf_docs = st.sidebar.file_uploader("Select Files", accept_multiple_files=True)
    # file_path=""
    # if pdf_docs:
    #     file_type+=save_uploadedfile(pdf_docs[0])

    # # Process PDF
    # if st.sidebar.button("Process PDF"):
    #     if pdf_docs:
    #         st.sidebar.write('PDF is processed')

    # Sidebar for YouTube URL input
    st.sidebar.header("Enter URL for (YOUTUBE/ WEBSITE)")
    youtube_url = st.sidebar.text_input("Enter URL")

    # # Process URL button
    # if st.sidebar.button("Process URL"):
    #     if youtube_url:
    #         st.sidebar.write("Processing URL:", youtube_url)
    #         if is_valid_youtube_url(youtube_url):
    #             st.sidebar.write('URL is valid')
    #         else:
    #             st.sidebar.write('URL is invalid')
    
    question = st.text_input('Enter your query')
    if st.button("Ask Question"):
        if question:
            # documents=""
            documents = load_documents(youtube_url,pdf_docs,file_type)
            split = split_documents(documents)
            db = define_vectore_store(split)
            rag_chain = retriever_chain(db)

            answer = rag_chain.invoke(question)
            idx = answer.find('Question:')
            answer = answer[idx:]
            st.write(answer)
        else:
            st.warning("Please enter a question")

if __name__=='__main__':
    main()